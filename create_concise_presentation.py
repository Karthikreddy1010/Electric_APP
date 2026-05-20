import os
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_concise_presentation():
    prs = Presentation()
    
    # Theme Colors
    DARK_BLUE = RGBColor(15, 32, 67)
    PRIMARY_BLUE = RGBColor(37, 99, 235)
    LIGHT_GRAY = RGBColor(241, 245, 249)

    # Slide Layouts
    TITLE_SLIDE_LAYOUT = prs.slide_layouts[0]
    BULLET_SLIDE_LAYOUT = prs.slide_layouts[1]
    
    def apply_title_style(title_shape):
        if not title_shape: return
        title_shape.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        
    def add_slide_with_image(title, bullets, image_name, img_left, img_top, img_width):
        slide = prs.slides.add_slide(BULLET_SLIDE_LAYOUT)
        title_shape = slide.shapes.title
        title_shape.text = title
        apply_title_style(title_shape)
        
        # Add bullets on the left side
        body_shape = slide.shapes.placeholders[1]
        body_shape.left = Inches(0.5)
        body_shape.width = Inches(4.0)
        tf = body_shape.text_frame
        tf.clear() 
        
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(22)
            p.font.color.rgb = DARK_BLUE
            p.space_after = Pt(14)
            
        # Add image on the right side (or placeholder)
        if os.path.exists(image_name):
            slide.shapes.add_picture(image_name, Inches(img_left), Inches(img_top), width=Inches(img_width))
        else:
            # Placeholder box
            shape = slide.shapes.add_shape(
                1, # MSO_SHAPE.RECTANGLE
                Inches(img_left), Inches(img_top), Inches(img_width), Inches(4.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = LIGHT_GRAY
            shape.line.color.rgb = PRIMARY_BLUE
            shape.line.width = Pt(2)
            ptf = shape.text_frame
            ptf.text = f"[Insert {image_name} here]"
            for paragraph in ptf.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.color.rgb = PRIMARY_BLUE
                paragraph.font.bold = True
                paragraph.font.size = Pt(16)
        
        return slide

    # Slide 1: Overview & Analytics
    add_slide_with_image(
        "1. Executive Overview & Analytics",
        [
            "Analyze actual electricity costs and track bill component breakdowns.",
            "Visualize data through interactive charts and graphical dashboards.",
            "Decompose delivery, generation, and fixed customer charges."
        ],
        "overview.png",
        img_left=5.0, img_top=1.5, img_width=4.5
    )

    # Slide 2: Trends & Forecasting
    add_slide_with_image(
        "2. Price Trends & Forecasting",
        [
            "Identify historical trends in electricity prices over time.",
            "Forecast future electricity demand using Prophet + SARIMA ensemble models.",
            "Track month-over-month (MoM) and year-over-year (YoY) volatility."
        ],
        "forecast.png",
        img_left=5.0, img_top=1.5, img_width=4.5
    )

    # Slide 3: Regional Benchmarking
    add_slide_with_image(
        "3. Regional Price Benchmarking",
        [
            "Benchmark electricity prices across the U.S. against national averages.",
            "Analyze specific New Jersey pricing relative to neighboring states.",
            "Identify macro-level geographic anomalies in energy rates."
        ],
        "benchmark.png",
        img_left=5.0, img_top=1.5, img_width=4.5
    )

    # Slide 4: Retail Supply Plans
    add_slide_with_image(
        "4. Evaluating Retail Supply Plans",
        [
            "Evaluate available retail supply plans for cost optimization.",
            "Compare fixed vs. variable charges and term durations.",
            "Simulate supply changes to predict projected bill impacts seamlessly."
        ],
        "plans.png",
        img_left=5.0, img_top=1.5, img_width=4.5
    )

    prs.save("ElectricAI_Visual_Deck.pptx")
    print("Presentation saved to ElectricAI_Visual_Deck.pptx")

if __name__ == "__main__":
    create_concise_presentation()
