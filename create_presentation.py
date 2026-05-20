import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Theme Colors
    DARK_BLUE = RGBColor(15, 32, 67)
    PRIMARY_BLUE = RGBColor(37, 99, 235)
    TEAL_GREEN = RGBColor(13, 148, 136)
    LIGHT_GRAY = RGBColor(241, 245, 249)

    # Slide Layouts
    TITLE_SLIDE_LAYOUT = prs.slide_layouts[0]
    BULLET_SLIDE_LAYOUT = prs.slide_layouts[1]
    
    # Helper functions
    def apply_title_style(title_shape):
        if not title_shape: return
        title_shape.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.size = Pt(40)
        
    def add_slide_with_bullets(title, bullets):
        slide = prs.slides.add_slide(BULLET_SLIDE_LAYOUT)
        title_shape = slide.shapes.title
        title_shape.text = title
        apply_title_style(title_shape)
        
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear() # Clear default
        
        for i, bullet in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            if bullet.startswith("  - "):
                p.text = bullet[4:]
                p.level = 1
                p.font.size = Pt(20)
                p.font.color.rgb = PRIMARY_BLUE
            elif bullet.startswith("- "):
                p.text = bullet[2:]
                p.level = 0
                p.font.size = Pt(24)
                p.font.color.rgb = DARK_BLUE
            else:
                p.text = bullet
                p.level = 0
                p.font.size = Pt(24)
                p.font.color.rgb = DARK_BLUE
        return slide

    def add_placeholder_box(slide, left, top, width, height, text):
        shape = slide.shapes.add_shape(
            1, # MSO_SHAPE.RECTANGLE
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_GRAY
        shape.line.color.rgb = PRIMARY_BLUE
        shape.line.width = Pt(2)
        tf = shape.text_frame
        tf.text = text
        for paragraph in tf.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.color.rgb = PRIMARY_BLUE
            paragraph.font.bold = True
            paragraph.font.size = Pt(16)

    # 1. Title Slide
    slide = prs.slides.add_slide(TITLE_SLIDE_LAYOUT)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "AI-Powered Energy Intelligence Platform"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "Explainable Billing, Forecasting, and Grid Analytics"
    subtitle.text_frame.paragraphs[0].font.color.rgb = TEAL_GREEN
    subtitle.text_frame.paragraphs[0].font.italic = True
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)

    # 2. Problem Statement
    add_slide_with_bullets(
        "The Problem with Energy Analytics",
        [
            "- Electricity billing is complex, opaque, and difficult to interpret",
            "- Customers lack transparency into bill fluctuations",
            "- Utilities struggle with demand forecasting and pricing strategy",
            "- Regional pricing and rate tracking are fragmented"
        ]
    )

    # 3. Objectives
    add_slide_with_bullets(
        "Platform Objectives",
        [
            "- Improve billing transparency using AI",
            "- Enable accurate demand forecasting at grid level",
            "- Quantify impact of pricing and usage changes",
            "- Provide regional and geographic intelligence",
            "- Support data-driven decision making"
        ]
    )

    # 4. Solution Overview
    slide = add_slide_with_bullets(
        "Solution Overview: System Architecture",
        [
            "- Show modules: Overview, Forecast, Impact, Benchmark, Geo",
            "- Data sources: PJM, EIA, weather, rate histories",
            "- AI/ML layer + visualization/dashboard layer"
        ]
    )
    add_placeholder_box(slide, 5, 2, 4.5, 4.5, "[Insert System Architecture Diagram Here]")

    # 5. Explainable Billing
    slide = add_slide_with_bullets(
        "Explainable Billing (Overview Module)",
        [
            "- Break down bills into components: delivery, generation, taxes, credits",
            "- Use LLMs to generate natural language explanations",
            "- Highlight anomaly detection (usage spikes, seasonal changes)"
        ]
    )
    add_placeholder_box(slide, 5, 1.5, 4.5, 5, "[Insert overview.png and bill_narrative.png Here]")

    # 6. Predictive Grid Modeling
    slide = add_slide_with_bullets(
        "Predictive Grid Modeling (Forecast Module)",
        [
            "- Use real-world PJM datasets (130MB+ scale)",
            "- Models: Prophet + SARIMA ensemble",
            "- Outputs: daily load forecasting, peak demand prediction"
        ]
    )
    add_placeholder_box(slide, 5, 2, 4.5, 4, "[Insert forecast.png Chart Here]")

    # 7. Deterministic & Causal Sensitivity
    slide = add_slide_with_bullets(
        "Deterministic & Causal Sensitivity",
        [
            "- Dual engine approach:",
            "  - A. Accounting Decomposition Engine (usage vs rate)",
            "  - B. Causal Inference Engine (HDD/CDD, DML)",
            "- Estimate true demand elasticity"
        ]
    )
    add_placeholder_box(slide, 5.5, 1.5, 4, 5, "[Insert impact_breakdown.png and simulator.png Here]")

    # 8. Regional Price Benchmarking
    slide = add_slide_with_bullets(
        "Regional Price Benchmarking",
        [
            "- Use EIA state-level pricing data",
            "- Compare rankings, volatility, and historical trends",
            "- Map and bar chart visual benchmarking"
        ]
    )
    add_placeholder_box(slide, 6, 2, 3.5, 4.5, "[Insert Regional Map/Bar Chart Here]")

    # 9. Geographic Rate Tracking
    slide = add_slide_with_bullets(
        "Geographic Rate Tracking (Geo Insights)",
        [
            "- ZIP code-level insights and aggregation",
            "- Track utility-specific rate histories (e.g., PSE&G)",
            "- Real-time alerts for rate changes"
        ]
    )
    add_placeholder_box(slide, 6, 2, 3.5, 4.5, "[Insert Geo Tab Map Here]")

    # 10. Key Capabilities Summary
    add_slide_with_bullets(
        "Key Capabilities Summary",
        [
            "- Explainability via large language models",
            "- High-accuracy forecasting (ensemble ML)",
            "- Advanced causal analysis (Double Machine Learning)",
            "- Geographic intelligence mapping",
            "- Interactive, real-time dashboards"
        ]
    )

    # 11. Tech Stack
    add_slide_with_bullets(
        "Technology Stack",
        [
            "- Backend: Python (FastAPI, Prophet, SARIMA, DoWhy)",
            "- Data Sources: PJM, EIA, NOAA Weather APIs",
            "- AI Integration: Ollama / Qwen local LLM",
            "- Visualization: React, Recharts, TailwindCSS"
        ]
    )

    # 12. Business Impact
    add_slide_with_bullets(
        "Business Impact & Value",
        [
            "- Increased transparency and trust for customers",
            "- Better demand planning and grid stability for utilities",
            "- Optimized pricing and margin strategies",
            "- Regulatory reporting and compliance support"
        ]
    )

    # 13. Future Enhancements
    add_slide_with_bullets(
        "Future Enhancements",
        [
            "- Real-time data streaming integration",
            "- Personalized, AI-driven user recommendations",
            "- Automated end-to-end reporting (PDF/LLM summaries)",
            "- Seamless integration with smart meter infrastructure"
        ]
    )

    # 14. Conclusion
    slide = prs.slides.add_slide(TITLE_SLIDE_LAYOUT)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Transforming Energy Analytics"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "A Scalable, Explainable, and Data-Driven Platform\n\nThank You"
    subtitle.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE

    # Save
    prs.save("ElectricAI_Investor_Presentation.pptx")
    print("Presentation saved successfully to ElectricAI_Investor_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
